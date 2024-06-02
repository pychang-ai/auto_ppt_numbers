import os
from pptx import Presentation
from pptx.util import Inches, Pt

def create_sample_pptx(pptx_path):
    prs = Presentation()
    slide_layout = prs.slide_layouts[1]  # 1表示标题和内容布局

    for i in range(1, 10):
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = f"Sample Slide {i}"
        content.text = f"This is the content of slide {i}."

    prs.save(pptx_path)
    print(f"Sample presentation created and saved as {pptx_path}")

def add_step_titles(pptx_path):
    if not os.path.exists(pptx_path):
        create_sample_pptx(pptx_path)

    prs = Presentation(pptx_path)
    total_slides = len(prs.slides)

    for idx, slide in enumerate(prs.slides):
        textbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(4), Inches(1))  # 位置和大小可以根据需要调整
        text_frame = textbox.text_frame
        p = text_frame.add_paragraph()
        p.text = f"Step {idx + 1}/{total_slides}"
        p.font.size = Pt(24)  # 字体大小可以根据需要调整

    updated_pptx_path = "updated_" + pptx_path
    prs.save(updated_pptx_path)
    print(f"Updated presentation saved as {updated_pptx_path}")

# 使用示例
if __name__ == "__main__":
    pptx_path = 'example_presentation.pptx'  # 请替换为你的 PowerPoint 文件路径
    add_step_titles(pptx_path)
